import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Color Palette
    navy = RGBColor(30, 58, 138)
    teal = RGBColor(15, 118, 110)
    charcoal = RGBColor(30, 41, 59)
    white = RGBColor(255, 255, 255)
    light_slate = RGBColor(241, 245, 249)
    amber = RGBColor(245, 158, 11)
    
    # Helper: Base slide with standard header
    def add_slide_base(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header block
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
        header.fill.solid()
        header.fill.fore_color.rgb = navy
        header.line.fill.background()
        
        # Title text
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = white
        
        # Divider line
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0), Inches(13.33), Inches(0.08))
        divider.fill.solid()
        divider.fill.fore_color.rgb = teal
        divider.line.fill.background()
        
        return slide

    # Helper: Add textbox bullets
    def add_bullets(slide, left, top, width, height, bullets, size=16):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(size)
            p.font.color.rgb = charcoal
            p.space_after = Pt(10)
            if bullet.startswith("• "):
                p.level = 0
            elif bullet.startswith("  - "):
                p.level = 1

    # Helper: Add picture with fallback placeholder
    def add_picture_safe(slide, image_path, left, top, width, height):
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        else:
            # Draw placeholder shape
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = light_slate
            rect.line.color.rgb = RGBColor(203, 213, 225)
            tb = slide.shapes.add_textbox(left, top, width, height)
            p = tb.text_frame.paragraphs[0]
            p.text = f"Screenshot / Chart:\n{os.path.basename(image_path)}\n(Run full pipeline to generate)"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = charcoal
            p.alignment = PP_ALIGN.CENTER

    # ── SLIDE 1: TITLE SLIDE ────────────────────────────────────────────────
    slide_1 = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = navy
    rect.line.fill.background()
    
    # Title textbox
    tb_1 = slide_1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2.5))
    tf_1 = tb_1.text_frame
    p_1 = tf_1.paragraphs[0]
    p_1.text = "Mutual Fund Analytics Platform"
    p_1.font.size = Pt(46)
    p_1.font.bold = True
    p_1.font.color.rgb = white
    p_1.alignment = PP_ALIGN.CENTER
    
    p_2 = tf_1.add_paragraph()
    p_2.text = "End-to-End ETL Pipeline, Performance Analytics & Power BI Dashboard"
    p_2.font.size = Pt(20)
    p_2.font.color.rgb = light_slate
    p_2.alignment = PP_ALIGN.CENTER
    p_2.space_before = Pt(15)
    
    # Details textbox
    tb_1_sub = slide_1.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.33), Inches(1.5))
    p_3 = tb_1_sub.text_frame.paragraphs[0]
    p_3.text = "Intern Name: Vikas Prajapat   |   Organization: Bluestock FinTech"
    p_3.font.size = Pt(14)
    p_3.font.color.rgb = RGBColor(226, 232, 240)
    p_3.alignment = PP_ALIGN.CENTER

    # ── SLIDE 2: PROBLEM STATEMENT ──────────────────────────────────────────
    slide_2 = add_slide_base("Problem Statement")
    bullets_2 = [
        "• Mutual Fund Complexity: Investors are overwhelmed by the volume and variety of mutual fund schemes, making manual comparison ineffective.",
        "• Invisible Risk Ratios: Traditional platforms report raw historical returns but fail to highlights risk metrics (Sharpe, Sortino, VaR, Max Drawdown).",
        "• Fragmented Data: NAV histories, investor demographics, and sector holdings reside in separate silos, preventing consolidated views.",
        "• High Cost Drag: Cost differences between Regular and Direct plans are often ignored by retail clients, eroding long-term compounding wealth."
    ]
    add_bullets(slide_2, Inches(0.75), Inches(1.8), Inches(6.5), Inches(5.0), bullets_2, size=18)
    
    # Highlight Box on the right
    box_2 = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5))
    box_2.fill.solid()
    box_2.fill.fore_color.rgb = light_slate
    box_2.line.color.rgb = teal
    tb_box2 = slide_2.shapes.add_textbox(Inches(8.2), Inches(2.2), Inches(4.1), Inches(3.7))
    tf_box2 = tb_box2.text_frame
    tf_box2.word_wrap = True
    p_b2 = tf_box2.paragraphs[0]
    p_b2.text = "The Challenge:"
    p_b2.font.size = Pt(22)
    p_b2.font.bold = True
    p_b2.font.color.rgb = navy
    p_b2.space_after = Pt(15)
    p_b2_sub = tf_box2.add_paragraph()
    p_b2_sub.text = "How can we build an integrated analytics pipeline to ingest raw AMC sheets, resolve missing holiday data, compute advanced risk metrics, and serve a dynamic dashboard to investors?"
    p_b2_sub.font.size = Pt(18)
    p_b2_sub.font.color.rgb = charcoal

    # ── SLIDE 3: PROJECT OBJECTIVES ─────────────────────────────────────────
    slide_3 = add_slide_base("Project Objectives")
    bullets_3 = [
        "• Automated Data Ingestion: Setup a Python scanner to ingest, clean, and validate 10 raw datasets representing AMC holdings, NAVs, and transactions.",
        "• Database Warehouse: Design a Star Schema SQLite DB to house clean Dimension (fund, date) and Fact (NAV, transactions, AUM, performance) tables.",
        "• Performance Analytics Scorecard: Programmatically compute CAGR (1y, 3y, 5y), Sharpe, Sortino, Alpha, Beta, and Max Drawdown.",
        "• Fund Recommender: Build an algorithmic engine that filters schemes by risk category and recommends the top 3 options based on Sharpe ratios.",
        "• BI Interactive Dashboard: Publish a 4-page Power BI dashboard detailing Industry Trends, Fund Performance, Investor Analytics, and Market Trends."
    ]
    add_bullets(slide_3, Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0), bullets_3, size=18)

    # ── SLIDE 4: DATA SOURCES OVERVIEW ──────────────────────────────────────
    slide_4 = add_slide_base("Data Sources & Database Schema")
    bullets_4_left = [
        "• Core Dimension Tables:",
        "  - dim_fund: Master scheme information, launch dates, categories.",
        "  - dim_date: Generated date dimension with day, month, quarter, and weekend flags.",
        "• Cleaned Data Volume:",
        "  - 37,000+ historical NAV records in fact_nav.",
        "  - 15,000+ retail transactions in fact_transactions.",
        "  - 40 major schemes tracked with performance ratios."
    ]
    bullets_4_right = [
        "• Core Fact Tables:",
        "  - fact_nav: Daily NAV records mapped to calendar dates.",
        "  - fact_transactions: Purchases, redemptions, age, income, state.",
        "  - fact_performance: Annualized returns, risk ratios, ratings.",
        "  - fact_aum: AMC level AUM and folio counts over quarters.",
        "• Schema Relationship Constraints:",
        "  - Enforced SQLite primary keys & foreign key linkages."
    ]
    add_bullets(slide_4, Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0), bullets_4_left, size=16)
    add_bullets(slide_4, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), bullets_4_right, size=16)

    # ── SLIDE 5: PROJECT ARCHITECTURE ───────────────────────────────────────
    slide_5 = add_slide_base("Project Architecture Workflow")
    bullets_5 = [
        "• ETL Phase: Read raw CSVs, request live NAV updates from AMFI API.",
        "• Cleaning: Duplicate purging, outlier limits, and forward-filling holidays.",
        "• Storage: Load facts/dims to SQLite database.",
        "• Analytics: Run computations for Sharpe, Sortino, Alpha, Beta, VaR, CVaR, Rolling Sharpe, and HHI concentrations.",
        "• Dashboard: Direct SQLite connector displays interactive page visuals.",
        "• Reporting: Programmatic PDF and PPTX deck builders package deliverables."
    ]
    add_bullets(slide_5, Inches(0.75), Inches(1.6), Inches(5.5), Inches(5.2), bullets_5, size=16)
    add_picture_safe(slide_5, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\project_architecture.png", Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8))

    # ── SLIDE 6: EDA HIGHLIGHTS 1 — MARKET TRENDS ───────────────────────────
    slide_6 = add_slide_base("EDA Highlights — Market Trends")
    # Placing NAV trend and AUM Growth charts side-by-side
    add_picture_safe(slide_6, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\nav_trend.png", Inches(0.75), Inches(1.6), Inches(5.5), Inches(3.2))
    add_picture_safe(slide_6, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\aum_growth.png", Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.2))
    
    bullets_6 = [
        "• NAV Growth: Large-cap funds demonstrate robust historical trajectories, with SBI and Nippon Bluechips exhibiting consistent returns.",
        "• AMC Dominance: AUM analysis indicates that SBI, ICICI, and HDFC represent over 50% of the industry's total AUM, highlighting high retail trust."
    ]
    add_bullets(slide_6, Inches(0.75), Inches(5.0), Inches(11.83), Inches(2.0), bullets_6, size=16)

    # ── SLIDE 7: EDA HIGHLIGHTS 2 — DEMOGRAPHICS & GEOGRAPHICS ──────────────
    slide_7 = add_slide_base("EDA Highlights — Demographics & Geography")
    add_picture_safe(slide_7, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\age_distribution.png", Inches(0.75), Inches(1.6), Inches(5.5), Inches(3.2))
    add_picture_safe(slide_7, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\state_distribution.png", Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.2))
    
    bullets_7 = [
        "• Key Age Group: The 25-40 cohort represents over 60% of all transaction volumes, driven by digital payment modes.",
        "• Geographic Concentration: Maharashtra, Karnataka, and Delhi contribute the highest investment volumes, highlighting urban concentration."
    ]
    add_bullets(slide_7, Inches(0.75), Inches(5.0), Inches(11.83), Inches(2.0), bullets_7, size=16)

    # ── SLIDE 8: PERFORMANCE ANALYTICS & SCORECARD ─────────────────────────
    slide_8 = add_slide_base("Performance Analytics & Scorecard")
    bullets_8 = [
        "• Ratios Calculation:",
        "  - Sharpe Ratio: Excess return relative to risk-free rate ($R_f=6.5\%$) per unit of volatility.",
        "  - Sortino Ratio: Penalizes only downside deviation.",
        "  - Alpha & Beta: Run OLS regressions against Nifty 100 benchmark index.",
        "• 0-100 Scorecard Ranking Formula:",
        "  - 3Yr Returns (30%) + Sharpe (25%) + Alpha (20%) + Expense (15% inverse) + Max Drawdown (10% inverse).",
        "  - Nippon Large Cap & SBI Bluechip lead scorecard ranks."
    ]
    add_bullets(slide_8, Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.2), bullets_8, size=16)
    add_picture_safe(slide_8, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\fund_scorecard_chart.png", Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8))

    # ── SLIDE 9: ADVANCED RISK ANALYTICS ────────────────────────────────────
    slide_9 = add_slide_base("Advanced Risk & Sector Concentration")
    bullets_9 = [
        "• Tail Risk: Value at Risk (95% VaR) and CVaR measures single-day loss expectations, showing Mid-Cap funds carry double the tail-risk of Bluechip funds.",
        "• Rolling Sharpe Ratios: Shows dynamic changes in risk-adjusted performance over 90-day rolling windows.",
        "• Sector HHI Concentration: Sum of squared sector weights. Axis Bluechip has low HHI (~1,500), showing high diversification."
    ]
    add_bullets(slide_9, Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.2), bullets_9, size=16)
    add_picture_safe(slide_9, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\reports\rolling_sharpe_chart.png", Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8))

    # ── SLIDE 10: DASHBOARD OVERVIEW — MARKET & FUND PERFORMANCE ───────────
    slide_10 = add_slide_base("Dashboard: Market & Performance Overview")
    add_picture_safe(slide_10, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\dashboard\screenshots\Industry_Overview.png", Inches(0.75), Inches(1.6), Inches(5.5), Inches(3.2))
    add_picture_safe(slide_10, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\dashboard\screenshots\Fund_Performance.png", Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.2))
    
    bullets_10 = [
        "• Industry Overview Dashboard (Left): Treemaps of AMC market share, category inflows, and total AUM metrics.",
        "• Fund Performance Dashboard (Right): Scorecard matrix, 3y return vs Sharpe scatter plot, and expense toggles."
    ]
    add_bullets(slide_10, Inches(0.75), Inches(5.0), Inches(11.83), Inches(2.0), bullets_10, size=16)

    # ── SLIDE 11: DASHBOARD OVERVIEW — INVESTOR ANALYTICS ───────────────────
    slide_11 = add_slide_base("Dashboard: Investor & SIP Trends")
    add_picture_safe(slide_11, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\dashboard\screenshots\Investor_Analytics.png", Inches(0.75), Inches(1.6), Inches(5.5), Inches(3.2))
    add_picture_safe(slide_11, r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\dashboard\screenshots\SIP_Market_Trends.png", Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.2))
    
    bullets_11 = [
        "• Investor Analytics Dashboard (Left): Age demographics, income distributions, and geographical state heatmaps.",
        "• SIP & Market Trends Dashboard (Right): Monthly inflows, SIP continuity indexes, and correlation with NAV fluctuations."
    ]
    add_bullets(slide_11, Inches(0.75), Inches(5.0), Inches(11.83), Inches(2.0), bullets_11, size=16)

    # ── SLIDE 12: KEY FINDINGS & THANK YOU ──────────────────────────────────
    slide_12 = add_slide_base("Key Findings & Contact Information")
    bullets_12_left = [
        "• Key Performance Findings:",
        "  - Top Sharpe: SBI Bluechip Fund (Sharpe Ratio = 1.76)",
        "  - Top CAGR: Nippon India Large Cap (3Yr CAGR = 22.45%)",
        "  - Top Alpha: ICICI Prudential Bluechip (Alpha = 3.12%)",
        "  - Top AUM: SBI Mutual Fund (> 48,000 Crores)",
        "• Valuable Cohort: Q1 2025 vintage shows 82% retention."
    ]
    add_bullets(slide_12, Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.2), bullets_12_left, size=16)
    
    # Thank you box on the right
    box_12 = slide_12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8))
    box_12.fill.solid()
    box_12.fill.fore_color.rgb = navy
    box_12.line.color.rgb = teal
    tb_box12 = slide_12.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.4), Inches(4.0))
    tf_box12 = tb_box12.text_frame
    tf_box12.word_wrap = True
    
    p_b12 = tf_box12.paragraphs[0]
    p_b12.text = "Thank You!"
    p_b12.font.size = Pt(36)
    p_b12.font.bold = True
    p_b12.font.color.rgb = white
    p_b12.alignment = PP_ALIGN.CENTER
    p_b12.space_after = Pt(25)
    
    p_b12_sub1 = tf_box12.add_paragraph()
    p_b12_sub1.text = "Mutual Fund Analytics Capstone Project"
    p_b12_sub1.font.size = Pt(18)
    p_b12_sub1.font.color.rgb = light_slate
    p_b12_sub1.alignment = PP_ALIGN.CENTER
    p_b12_sub1.space_after = Pt(10)
    
    p_b12_sub2 = tf_box12.add_paragraph()
    p_b12_sub2.text = "Presenter: Vikas Prajapat\nContact: vikas.prajapat@example.com\nOrganization: Bluestock FinTech"
    p_b12_sub2.font.size = Pt(16)
    p_b12_sub2.font.color.rgb = RGBColor(226, 232, 240)
    p_b12_sub2.alignment = PP_ALIGN.CENTER
    
    # Save Presentation
    output_dir = r"d:\Mutual Fund Analytic\Mutual-Fund-Analytics\presentation"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "Bluestock_MF_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation deck successfully saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
